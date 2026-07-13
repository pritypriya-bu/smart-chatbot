"""
creators.py - Generate new files from chat instructions:

  - Excel (.xlsx)      - from existing data or a description
  - PDF (.pdf)         - text-based reports/documents
  - PowerPoint (.pptx) - slide decks from a topic/outline
  - Chart image (.png) - from a loaded DataFrame (matplotlib)
  - AI image (.png)    - from a text prompt (Pollinations.ai, free, no key)

Everything is free. AI images require internet access (Pollinations.ai).
"""

from __future__ import annotations
import io
import re
import urllib.parse
import requests
import pandas as pd


# ---------------------------------------------------------------------------
# INTENT DETECTION - is the user asking us to create a file?
# ---------------------------------------------------------------------------
_CREATE_WORDS = [
    "create", "make", "generate", "build", "banao", "bana do", "banade",
    "bana de", "banana", "export", "save as", "save to", "nikaalo", "tayar",
    "draft", "convert", "turn into", "turn this into", "draw", "plot",
    # follow-up phrasings ("another image", "one more chart")
    "another", "one more", "ek aur", "aur ek", "one more time",
]

_TYPE_PATTERNS = {
    "excel": ["excel", "xlsx", "spreadsheet", "sheet", "workbook"],
    "pptx":  ["presentation", "powerpoint", "ppt", "pptx", "slide", "slides", "deck"],
    "pdf":   ["pdf", "report"],
    "chart": ["chart", "graph", "plot", "bar chart", "line chart", "pie",
              "histogram", "scatter", "visualize", "visualization"],
    "image": ["image", "picture", "photo", "tasveer", "drawing", "illustration",
              "art", "logo", "poster", "wallpaper"],
}


def _has_word(text, phrase):
    """Whole-word/phrase match: 'art' matches 'art' but not 'part' or 'chart'."""
    return re.search(r"(?<!\w)" + re.escape(phrase) + r"(?!\w)", text) is not None


def detect_creation(prompt: str):
    """
    Return the target file type ("excel"/"pptx"/"pdf"/"chart"/"image") if
    the prompt asks to create a file, or None.
    """
    # Strip filenames (e.g. policy.pdf, data.xlsx) - they refer to a SOURCE
    # file, not the output format we should create.
    p = re.sub(r"\b[\w\-]+\.[a-z0-9]{2,5}\b", " ", prompt.lower())
    has_create = any(_has_word(p, w) for w in _CREATE_WORDS)

    # Without an explicit "create/make/..." word we don't trigger creation.
    if not has_create:
        return None

    # Match type in priority order; all whole-word.
    for t in ("chart", "pptx", "excel", "pdf", "image"):
        if any(_has_word(p, w) for w in _TYPE_PATTERNS[t]):
            return t
    return None


# ---------------------------------------------------------------------------
# EXCEL
# ---------------------------------------------------------------------------
def create_excel(llm, prompt, df=None):
    """Export the provided DataFrame, or ask the LLM to generate one."""
    if df is not None and _wants_existing_data(prompt):
        buf = io.BytesIO()
        with pd.ExcelWriter(buf, engine="openpyxl") as xl:
            df.to_excel(xl, index=False, sheet_name="Data")
        return (buf.getvalue(), "data.xlsx",
                f"Exported current data ({len(df)} rows) to Excel.")

    # LLM-generated table
    sys = ("You generate tabular data as JSON. Return ONLY JSON, no prose:\n"
           '{"sheet_name": "...", "columns": ["c1","c2",...], '
           '"rows": [["v1","v2",...], ...]}\n'
           "Make realistic, useful data matching the request.")
    plan = llm.ask_json(sys, prompt)
    if not isinstance(plan, dict) or "columns" not in plan:
        # Fallback: if a DataFrame is loaded, export that instead
        if df is not None:
            return create_excel(llm, "export current data", df)
        raise ValueError("Could not generate structured data for the Excel file.")

    out_df = pd.DataFrame(plan.get("rows", []), columns=plan["columns"])
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as xl:
        out_df.to_excel(buf if False else xl, index=False,
                        sheet_name=str(plan.get("sheet_name", "Sheet1"))[:31])
    return (buf.getvalue(), "generated.xlsx",
            f"Excel created: {len(out_df)} rows x {len(out_df.columns)} columns.")


def _wants_existing_data(prompt):
    """Heuristic: does the user want to export the currently loaded data?"""
    p = prompt.lower()
    return any(w in p for w in ["this data", "current", "loaded", "ye data",
                                "is data", "filtered", "view", "table ko",
                                "isko", "iss data"])


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def create_pdf(llm, prompt, context_text=None):
    """Create a document PDF. When context_text is given, base it on that."""
    sys = (
        "You are a writer producing a clean PDF. Return ONLY JSON, no prose.\n"
        "Schema:\n"
        '{"title": "...",\n'
        ' "sections": [\n'
        '   {"heading": "...", "body": "paragraph text (optional)",\n'
        '    "bullets": ["optional bullet", ...],\n'
        '    "table": {"columns": ["c1","c2"], "rows": [["v1","v2"], ...]}  (optional) }\n'
        ' ]}\n'
        "Rules:\n"
        "- Use 'table' whenever the source material is comparative / tabular "
        "  (e.g. multiple items with the same attributes like weather for several cities, "
        "  product specs, results per row). Pick concise columns.\n"
        "- Use 'bullets' for short lists of points.\n"
        "- Use 'body' for narrative paragraphs.\n"
        "- 3-8 sections total. Every section must have a heading."
    )
    user = prompt
    if context_text:
        ctx = context_text[:12000]
        user = f"{prompt}\n\nBase it on this source material:\n{ctx}"
    plan = llm.ask_json(sys, user)
    if not isinstance(plan, dict):
        plan = {"title": "Document",
                "sections": [{"heading": "", "body": str(plan or prompt)}]}

    return (_build_pdf(plan), "document.pdf",
            f"PDF created: '{plan.get('title', 'Document')}'.")


def _build_pdf(plan):
    """Render a plan dict into a PDF byte string."""
    from fpdf import FPDF

    def clean(s):
        # fpdf core fonts only support latin-1 - encode safely
        return str(s).encode("latin-1", "replace").decode("latin-1")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()

    pdf.set_font("Helvetica", "B", 20)
    pdf.multi_cell(0, 10, clean(plan.get("title", "Document")))
    pdf.ln(4)

    for sec in plan.get("sections", []):
        pdf.set_x(pdf.l_margin)
        head = sec.get("heading", "")
        if head:
            pdf.set_font("Helvetica", "B", 14)
            pdf.multi_cell(0, 8, clean(head))
            pdf.ln(1)

        body = sec.get("body", "")
        if body:
            pdf.set_x(pdf.l_margin)
            pdf.set_font("Helvetica", "", 11)
            pdf.multi_cell(0, 6, clean(body))
            pdf.ln(2)

        bullets = sec.get("bullets") or []
        if bullets:
            pdf.set_font("Helvetica", "", 11)
            for b in bullets:
                pdf.set_x(pdf.l_margin)
                pdf.multi_cell(0, 6, clean(f"- {b}"))
            pdf.ln(2)

        table = sec.get("table") or {}
        cols = table.get("columns") or []
        rows = table.get("rows") or []
        if cols and rows:
            _draw_pdf_table(pdf, cols, rows, clean)
            pdf.ln(2)

        pdf.ln(1)

    out = pdf.output()
    return bytes(out)


def _draw_pdf_table(pdf, cols, rows, clean):
    """Draw a simple table with a header row inside the PDF."""
    page_w = pdf.w - pdf.l_margin - pdf.r_margin
    n = max(len(cols), 1)
    col_w = page_w / n
    line_h = 6

    pdf.set_x(pdf.l_margin)

    # Header
    pdf.set_font("Helvetica", "B", 10)
    pdf.set_fill_color(230, 230, 230)
    for c in cols:
        pdf.cell(col_w, line_h + 2, clean(str(c))[:60], border=1, fill=True)
    pdf.ln(line_h + 2)
    pdf.set_x(pdf.l_margin)

    # Rows
    pdf.set_font("Helvetica", "", 10)
    for r in rows:
        vals = list(r) + [""] * (n - len(r))
        for v in vals[:n]:
            pdf.cell(col_w, line_h + 2, clean(str(v))[:60], border=1)
        pdf.ln(line_h + 2)
        pdf.set_x(pdf.l_margin)


# ---------------------------------------------------------------------------
# POWERPOINT
# ---------------------------------------------------------------------------
def create_pptx(llm, prompt):
    """Generate a PowerPoint deck from a topic/outline."""
    sys = ("You create presentation outlines. Return ONLY JSON, no prose:\n"
           '{"title": "Deck Title", "subtitle": "optional", '
           '"slides": [{"title": "Slide title", "bullets": ["point 1","point 2"]}]}\n'
           "5-8 content slides, 3-5 concise bullets each.")
    plan = llm.ask_json(sys, prompt)
    if not isinstance(plan, dict) or "slides" not in plan:
        raise ValueError("Could not generate a presentation outline.")
    return (_build_pptx(plan), "presentation.pptx",
            f"Presentation created: {len(plan.get('slides', []))} slides.")


def _build_pptx(plan):
    """Render an outline dict into a .pptx byte string."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()

    # Title slide
    title_layout = prs.slide_layouts[0]
    s = prs.slides.add_slide(title_layout)
    s.shapes.title.text = str(plan.get("title", "Presentation"))
    if s.placeholders and len(s.placeholders) > 1 and plan.get("subtitle"):
        s.placeholders[1].text = str(plan["subtitle"])

    # Content slides
    bullet_layout = prs.slide_layouts[1]
    for slide in plan.get("slides", []):
        sl = prs.slides.add_slide(bullet_layout)
        sl.shapes.title.text = str(slide.get("title", ""))
        body = sl.placeholders[1].text_frame
        body.clear()
        bullets = slide.get("bullets", []) or [""]
        for i, b in enumerate(bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = str(b)
            para.font.size = Pt(18)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# CHART (data -> PNG via matplotlib)
# ---------------------------------------------------------------------------
def create_chart(llm, prompt, df):
    """Pick a chart spec from the LLM and render it against the DataFrame."""
    if df is None:
        raise ValueError("Load a data file first before creating a chart.")
    sys = ("You pick a chart spec for a pandas DataFrame. Return ONLY JSON:\n"
           '{"kind": "bar|line|pie|scatter|hist", "x": "col", "y": "col or null", '
           '"agg": "sum|mean|count|none", "title": "..."}\n'
           f"Available columns: {', '.join(map(str, df.columns))}")
    spec = llm.ask_json(sys, prompt) or {}
    return (_build_chart(df, spec, prompt), "chart.png",
            f"Chart created ({spec.get('kind', 'bar')}).")


def _resolve(col, df):
    """Case-insensitive column-name lookup."""
    if col in df.columns:
        return col
    low = {str(c).lower(): c for c in df.columns}
    return low.get(str(col).lower()) if col else None


def _build_chart(df, spec, prompt):
    """Render a chart spec to PNG bytes."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    kind = (spec.get("kind") or "bar").lower()
    x = _resolve(spec.get("x"), df)
    y = _resolve(spec.get("y"), df)
    agg = (spec.get("agg") or "none").lower()
    title = spec.get("title") or "Chart"

    fig, ax = plt.subplots(figsize=(8, 5))

    try:
        if kind == "hist" and x:
            df[x].plot(kind="hist", ax=ax, bins=20)
        elif kind == "pie" and x:
            counts = df[x].value_counts().head(10)
            counts.plot(kind="pie", ax=ax, autopct="%1.0f%%")
            ax.set_ylabel("")
        elif kind == "scatter" and x and y:
            df.plot(kind="scatter", x=x, y=y, ax=ax)
        elif x and y:
            data = df
            if agg in ("sum", "mean", "count"):
                grp = df.groupby(x)[y]
                data = getattr(grp, agg)().reset_index()
            data.plot(kind=kind if kind in ("bar", "line") else "bar",
                      x=x, y=y, ax=ax, legend=True)
        elif x:
            df[x].value_counts().head(15).plot(kind="bar", ax=ax)
        else:
            df.select_dtypes("number").plot(kind="bar", ax=ax)
    except Exception:
        # Very safe fallback
        df.select_dtypes("number").head(20).plot(ax=ax)

    ax.set_title(title)
    fig.tight_layout()
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=130)
    plt.close(fig)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# AI IMAGE (Pollinations.ai - free, no API key; internet required)
# ---------------------------------------------------------------------------
_IMG_PROMPT_SYS = (
    "You turn a user's request into a short, vivid image-generation prompt "
    "(English, one line, <200 chars). If the user says things like 'another', "
    "'one more', 'same', 'again', 'different one', use the PREVIOUS image "
    "described in the conversation as the base subject. Never return an empty "
    "prompt. Return ONLY the prompt text - no quotes, no JSON, no explanation."
)


def _refine_image_prompt(llm, user_msg, context_text=None):
    """Use the LLM to turn the conversation into a good image prompt."""
    ctx = (context_text or "").strip()
    if ctx:
        user = (f"Conversation so far:\n{ctx}\n\n"
                f"Latest user message: {user_msg}\n\nImage prompt:")
    else:
        user = f"User message: {user_msg}\n\nImage prompt:"
    try:
        out = llm.ask(_IMG_PROMPT_SYS, user, temperature=0.4)
        out = (out or "").strip().strip('"').strip("'").splitlines()[0][:220]
        return out or user_msg
    except Exception:
        return user_msg


def create_image(prompt, width=1024, height=1024, subject_override=None):
    """Generate an AI image via Pollinations.ai (free, no API key)."""
    if subject_override:
        subject = subject_override.strip()
    else:
        # Strip common creation/filler words from the prompt so the remainder
        # is a compact subject description.
        subject = re.sub(
            r"(?i)\b(create|make|generate|banao|bana do|an?|the|image|picture|"
            r"photo|of|please|ki|ka|ek|tasveer|another|one more|again)\b",
            " ", prompt,
        )
        subject = re.sub(r"\s+", " ", subject).strip() or prompt
    enc = urllib.parse.quote(subject)
    url = (f"https://image.pollinations.ai/prompt/{enc}"
           f"?width={width}&height={height}&nologo=true")
    try:
        r = requests.get(url, timeout=120)
        r.raise_for_status()
    except requests.exceptions.RequestException as e:
        raise RuntimeError(
            f"Could not generate the image (internet/service issue): {e}"
        )
    if not r.headers.get("content-type", "").startswith("image"):
        raise RuntimeError("The image service did not return an image. Please try again later.")
    return r.content, "image.png", f"AI image created: \"{subject}\""
